from pathlib import Path
from typing import Any


def read_uploaded_file_with_diagnostics(uploaded_file) -> dict[str, Any]:
    if uploaded_file is None:
        return _reader_result("", "", "", {})

    uploaded_file.seek(0)
    file_name = uploaded_file.name
    suffix = Path(file_name).suffix.lower()

    if suffix == ".txt":
        text = uploaded_file.read().decode("utf-8", errors="ignore")
        cleaned_text = _clean_extracted_lines(text.splitlines())
        return _reader_result(file_name, "txt", cleaned_text, {})

    if suffix == ".pdf":
        return _reader_result(file_name, "pdf", _read_pdf(uploaded_file), {})

    if suffix == ".docx":
        return _read_docx(uploaded_file, file_name)

    if suffix == ".pptx":
        return _reader_result(file_name, "pptx", _read_pptx(uploaded_file), {})

    return _reader_result(file_name, suffix.lstrip("."), "", {})


def read_uploaded_file(uploaded_file) -> str:
    return read_uploaded_file_with_diagnostics(uploaded_file)["text"]


def _read_pdf(uploaded_file) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as error:
        raise RuntimeError("Install pypdf to read PDF files.") from error

    reader = PdfReader(uploaded_file)
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(page for page in pages if page.strip())


def _read_docx(uploaded_file, file_name: str) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError as error:
        raise RuntimeError("Install python-docx to read DOCX files.") from error

    document = Document(uploaded_file)
    text_parts = []

    for paragraph in document.paragraphs:
        text_parts.append(paragraph.text)

    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                text_parts.append(cell.text)

    cleaned_text = _clean_extracted_lines(text_parts)
    diagnostics = {
        "file_name": file_name,
        "paragraph_count": len(document.paragraphs),
        "table_count": len(document.tables),
        "extracted_character_count": len(cleaned_text),
    }

    return _reader_result(file_name, "docx", cleaned_text, diagnostics)


def _read_pptx(uploaded_file) -> str:
    try:
        from pptx import Presentation
    except ImportError as error:
        raise RuntimeError("Install python-pptx to read PPTX files.") from error

    presentation = Presentation(uploaded_file)
    text_parts = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())

    return "\n".join(text_parts)


def _clean_extracted_lines(lines: list[str]) -> str:
    cleaned_lines = []

    for line in lines:
        cleaned_line = " ".join((line or "").split())
        if cleaned_line:
            cleaned_lines.append(cleaned_line)

    return "\n".join(cleaned_lines)


def _reader_result(
    file_name: str,
    file_type: str,
    text: str,
    diagnostics: dict[str, Any],
) -> dict[str, Any]:
    return {
        "file_name": file_name,
        "file_type": file_type,
        "text": text,
        "diagnostics": diagnostics,
    }
