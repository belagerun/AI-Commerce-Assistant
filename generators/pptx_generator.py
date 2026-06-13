import re
from datetime import datetime
from pathlib import Path

from config.settings import ARTIFACTS_DIR


def generate_pptx_presentation(title: str, content: str, sources: list[str] | None = None) -> dict[str, str | int]:
    try:
        from pptx import Presentation
    except ImportError as error:
        raise RuntimeError("Install python-pptx to generate PPTX presentations.") from error

    ARTIFACTS_DIR.mkdir(parents=True, exist_ok=True)
    output_path = ARTIFACTS_DIR / f"ecommerce_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title or "E-commerce Customer Service Brief"
    title_slide.placeholders[1].text = "AI Agent Network"

    for index, section in enumerate(_sections(content), start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = section["title"] or f"Point {index}"
        body = slide.shapes.placeholders[1].text_frame
        body.clear()

        for bullet_index, bullet in enumerate(section["bullets"]):
            paragraph = body.paragraphs[0] if bullet_index == 0 else body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0

    if sources:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Store Materials Used"
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for index, source in enumerate(sources[:6]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = source

    presentation.save(output_path)
    return _artifact_result(output_path, "pptx")


def _sections(content: str) -> list[dict[str, list[str] | str]]:
    sentences = [
        sentence.strip()
        for sentence in re.split(r"(?<=[.!?])\s+", " ".join((content or "").split()))
        if sentence.strip()
    ]
    sentences = sentences[:24] or ["No content available."]
    sections = []

    for index in range(0, len(sentences), 4):
        group = sentences[index:index + 4]
        sections.append(
            {
                "title": f"Key Point {len(sections) + 1}",
                "bullets": [item[:170] for item in group],
            }
        )

        if len(sections) >= 6:
            break

    return sections


def _artifact_result(path: Path, artifact_type: str) -> dict[str, str | int]:
    return {
        "file_name": path.name,
        "file_path": str(path),
        "artifact_type": artifact_type,
        "file_size": path.stat().st_size,
    }
